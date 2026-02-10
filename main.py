import sys
import streamlit as st
import google.generativeai as genai
import os
from dotenv import load_dotenv
import fitz  # PyMuPDF
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import io
import json
from docx import Document  # لقراءة ملفات Word
from pptx import Presentation  # لقراءة ملفات PowerPoint

# 1. تحميل الإعدادات
load_dotenv()
st.set_page_config(page_title="Flomind Quiz Generator", page_icon="🐥", layout="wide")

# الكود الذكي لتحديد مسار Tesseract
if sys.platform.startswith('win'):
    # هذا المسار خاص بجهازك المحلي (ويندوز)
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
else:
    # في سيرفرات لينكس (GitHub/Streamlit Cloud)، هو يعرف المسار تلقائياً
    print("Assuming Linux environment for Tesseract")

# إعداد Google API
api_key = os.getenv("GOOGLE_API_KEY")
if api_key:
    genai.configure(api_key=api_key)

# دوال تحسين الصور للـ OCR
def preprocess_image_for_ocr(image):
    """
   تنظيف الصورة قبل قرائتها لتحسين النتائج
    """
    # 1. تحويل للصورة الرمادية (Grayscale)
    image = image.convert('L')
    
    # 2. زيادة التباين (Contrast) لتوضيح الحروف الباهتة
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(2.0)  # ضاعفنا التباين
    
    # 3. زيادة الحدة (Sharpness)
    enhancer = ImageEnhance.Sharpness(image)
    image = enhancer.enhance(1.5)

    return image

# --- دوال استخراج النصوص حسب نوع الملف ---
def extract_text_from_docx(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_txt(file):
    return file.read().decode("utf-8")

def extract_text_from_pdf(file):
    text = ""
    try:
        pdf_document = fitz.open(stream=file.read(), filetype="pdf")
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            page_text = page.get_text()
            
            # إذا كان النص قليلاً (صورة)، شغل الـ OCR المحسن
            if len(page_text.strip()) < 50:
                pix = page.get_pixmap(dpi=300) # دقة عالية
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                
                # تطبيق التحسين قبل القراءة باستخدام OCR
                processed_img = preprocess_image_for_ocr(img)
                
                # إعدادات خاصة لقراءة كتل النصوص (psm 6)
                custom_config = r'--oem 3 --psm 6' 
                ocr_text = pytesseract.image_to_string(processed_img, lang='ara+eng', config=custom_config)
                text += ocr_text + "\n"
            else:
                text += page_text + "\n"
    except Exception as e:
        st.error(f"Error in PDF extraction: {e}")
    return text

# دالة اختيار أفضل موديل متاح
def get_best_model():
    # نحاول استخدام Flash لأنه الأفضل، وإذا لم نجد نستخدم Pro
    try:
        available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        for m in available_models:
            if 'flash' in m: return m
        return "models/gemini-1.5-flash"
    except:
        return "models/gemini-pro"

# دالة توليد الأسئلة من النص باستخدام الموديل المختار
def get_questions(text, number_of_questions=5):
    if not api_key: return None
    
    # إذا كان النص طويلاً جداً، نأخذ أول 10000 حرف لتجنب امتلاء الكوتا
    if len(text) > 10000:
        text = text[:10000] + "\n...(Text truncated due to quota limits)"

    model_name = get_best_model()
    generation_config = {"temperature": 0.3} # قللنا الحرارة لدقة أكثر
    
    model = genai.GenerativeModel(model_name=model_name, generation_config=generation_config)

    prompt = f"""
    Role: You are an expert Professor creating an exam.
    Task: Create {number_of_questions} multiple-choice questions (MCQ) based on the provided text.
    
    🔴 CRITICAL INSTRUCTIONS (Follow Strictly):
    1. **Context Repair:** The text comes from OCR and may contain typos (e.g., "LMéteyub" instead of "Interrupt"). You MUST infer the correct technical terms based on the context before asking.
    2. **Scope Constraint:** You are allowed to use your general knowledge to clarify concepts mentioned in the text, BUT you must NOT ask about topics completely absent from the text. Stick to the provided subject matter (e.g., if the text is about OS, don't ask about Networking unless linked).
    3. **Question Quality:** - Avoid trivial word-matching questions.
       - Ask conceptual questions that test understanding.
       - Distractors (wrong options) must be plausible and related to the field, not random or obviously wrong.
    4. **Language:** The questions must be in the same language as the majority of the text (Arabic or English).

    Output Format (Valid JSON Only):
    {{
      "questions": [
        {{
          "id": 1,
          "question": "Clear and precise question?",
          "options": ["Correct Answer", "Distractor 1", "Distractor 2", "Distractor 3"],
          "correct_answer": "Correct Answer",
          "explanation": "Brief explanation of why this is correct based on the text context."
        }}
      ]
    }}
    
    Source Text:
    '''{text}'''
    """
    try:
        response = model.generate_content(prompt)
        clean_json = response.text.replace('```json', '').replace('```', '').strip()
        return json.loads(clean_json)
    except Exception as e:
        st.error(f"AI Generation Error: {e}")
        return None


def main():
    st.title("🐥Flomind Quiz Generator")
    st.markdown("Turn your study materials into interactive quizzes 🪄")

    if "quiz_data" not in st.session_state:
        st.session_state.quiz_data = None

    if not api_key:
        st.warning("⚠️ الرجاء وضع API Key")
        return

    # 1. تحديث قائمة الملفات المدعومة
    uploaded_file = st.file_uploader("ارفع ملفك الدراسي", type=["pdf", "docx", "pptx", "txt"])
    num_q = st.number_input("عدد الأسئلة", 1, 10, 5)

    if st.button("🚀 إنشاء الاختبار", type="primary"):
        if uploaded_file:
            with st.spinner("جاري تحليل الملف واستخراج النصوص..."):
                file_ext = uploaded_file.name.split('.')[-1].lower()
                extracted_text = ""
                
                # إعادة المؤشر للبداية
                uploaded_file.seek(0)

                # نوع الملف المرفوع
                try:
                    if file_ext == "pdf":
                        extracted_text = extract_text_from_pdf(uploaded_file)
                    elif file_ext == "docx":
                        extracted_text = extract_text_from_docx(uploaded_file)
                    elif file_ext == "pptx":
                        extracted_text = extract_text_from_pptx(uploaded_file)
                    elif file_ext == "txt":
                        extracted_text = extract_text_from_txt(uploaded_file)
                except Exception as e:
                    st.error(f"فشلت قراءة الملف: {e}")

                # التأكد من وجود نص
                if extracted_text and len(extracted_text.strip()) > 10:
                    # عرض جزء من النص (Debugging)
                    #with st.expander("👀 معاينة النص المستخرج (تأكد من الجودة)"):
                    #   st.text(extracted_text[:1000])

                    # توليد الأسئلة
                    result = get_questions(extracted_text, num_q)
                    if result:
                        st.session_state.quiz_data = result.get("questions", [])
                        st.success("تم إنشاء الاختبار بنجاح! 🎉")
                else:
                    st.error("لم يتم العثور على نصوص قابلة للقراءة في الملف.")
        else:
            st.warning("الرجاء رفع ملف أولاً.")

    st.divider()

    # عرض الأسئلة
    if st.session_state.quiz_data:
        for q in st.session_state.quiz_data:
            st.subheader(f"{q['id']}. {q['question']}")
            user_choice = st.radio(f"Select:", q['options'], index=None, key=f"q_{q['id']}")
            
            if st.button(f"تأكيد الإجابة {q['id']}", key=f"btn_{q['id']}"):
                if user_choice == q['correct_answer']:
                    st.success("✅ إجابة صحيحة!")
                elif user_choice:
                    st.error(f"❌ خطأ. الصحيح: {q['correct_answer']}")
                st.info(f"💡 {q['explanation']}")
            st.markdown("---")
            
        if st.button("🔄 اختبار جديد"):
            st.session_state.quiz_data = None
            st.rerun()

if __name__ == '__main__':
    main()
